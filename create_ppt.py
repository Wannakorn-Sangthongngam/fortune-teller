from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color Palette (Deep Mystical Purple/Gold) ─────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x0D, 0x2B)   # deep navy/purple
BG_CARD      = RGBColor(0x1A, 0x1A, 0x3E)   # slightly lighter purple
ACCENT_GOLD  = RGBColor(0xF0, 0xC0, 0x40)   # gold
ACCENT_PURPLE= RGBColor(0x9B, 0x59, 0xB6)   # violet
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xCC, 0xCC, 0xFF)   # lavender white
TEXT_GOLD    = RGBColor(0xF0, 0xC0, 0x40)
DIVIDER      = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=BG_DARK):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color, alpha=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=18, bold=False, color=TEXT_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def divider(slide, t, color=DIVIDER, thickness=0.04):
    box(slide, 0.5, t, 12.33, thickness, color)


def section_tag(slide, label, l, t):
    """Small colored pill label."""
    box(slide, l, t, 1.8, 0.32, ACCENT_PURPLE)
    txt(slide, label, l + 0.05, t + 0.02, 1.7, 0.3,
        size=10, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)


def bullet_item(slide, text, l, t, w=10, size=16, color=TEXT_LIGHT, indent="  ▸  "):
    txt(slide, indent + text, l, t, w, 0.4, size=size, color=color)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s1 = add_slide()
bg(s1)

# Decorative top bar
box(s1, 0, 0, 13.33, 0.12, ACCENT_GOLD)

# Center glow box
box(s1, 2.5, 1.6, 8.33, 4.3, BG_CARD)
box(s1, 2.5, 1.6, 0.08, 4.3, ACCENT_GOLD)    # left accent bar

# Title
txt(s1, "🔮 Fortune Teller", 2.8, 2.0, 8.0, 1.2,
    size=48, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Subtitle
txt(s1, "AI-Powered Tarot Card Recognition", 2.8, 3.1, 8.0, 0.7,
    size=22, color=TEXT_LIGHT, align=PP_ALIGN.CENTER, italic=True)

divider(s1, 4.0, ACCENT_PURPLE, 0.03)

txt(s1, "EfficientNet-B0  ·  PyTorch  ·  Transfer Learning", 2.8, 4.1, 8.0, 0.5,
    size=14, color=TEXT_GOLD, align=PP_ALIGN.CENTER)

txt(s1, "Wannakorn Sangthongngam  ·  2026", 2.8, 4.8, 8.0, 0.4,
    size=13, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Bottom bar
box(s1, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Project Introduction
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s2 = add_slide()
bg(s2)
box(s2, 0, 0, 13.33, 0.12, ACCENT_GOLD)
box(s2, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)

section_tag(s2, "INTRODUCTION", 0.5, 0.25)
txt(s2, "What is Fortune Teller?", 0.5, 0.7, 10, 0.7,
    size=32, bold=True, color=ACCENT_GOLD)
divider(s2, 1.5)

# Two column boxes
box(s2, 0.5, 1.7, 5.7, 2.8, BG_CARD)
box(s2, 0.5, 1.7, 0.08, 2.8, ACCENT_GOLD)
txt(s2, "Overview", 0.75, 1.75, 5.2, 0.45, size=16, bold=True, color=ACCENT_GOLD)
bullet_item(s2, "Take a photo of a tarot card", 0.65, 2.2)
bullet_item(s2, "AI identifies which card it is", 0.65, 2.6)
bullet_item(s2, "Receive your fortune reading", 0.65, 3.0)
bullet_item(s2, "78 unique tarot cards supported", 0.65, 3.4)

box(s2, 6.8, 1.7, 5.9, 2.8, BG_CARD)
box(s2, 6.8, 1.7, 0.08, 2.8, ACCENT_PURPLE)
txt(s2, "Tech Stack", 7.05, 1.75, 5.4, 0.45, size=16, bold=True, color=ACCENT_PURPLE)
bullet_item(s2, "Framework: PyTorch", 6.95, 2.2)
bullet_item(s2, "Model: EfficientNet-B0", 6.95, 2.6)
bullet_item(s2, "Training: Google Colab (T4 GPU)", 6.95, 3.0)
bullet_item(s2, "Data: Custom tarot card photos", 6.95, 3.4)

# Flow diagram
box(s2, 1.0, 5.1, 2.2, 0.7, ACCENT_PURPLE)
txt(s2, "📷 Photo", 1.0, 5.2, 2.2, 0.5, size=16, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

txt(s2, "→", 3.3, 5.2, 0.6, 0.5, size=24, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

box(s2, 4.0, 5.1, 2.2, 0.7, ACCENT_PURPLE)
txt(s2, "🤖 AI Model", 4.0, 5.2, 2.2, 0.5, size=16, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

txt(s2, "→", 6.3, 5.2, 0.6, 0.5, size=24, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

box(s2, 7.0, 5.1, 2.2, 0.7, ACCENT_PURPLE)
txt(s2, "🃏 Card ID", 7.0, 5.2, 2.2, 0.5, size=16, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)

txt(s2, "→", 9.3, 5.2, 0.6, 0.5, size=24, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

box(s2, 10.0, 5.1, 2.5, 0.7, ACCENT_GOLD)
txt(s2, "🔮 Fortune", 10.0, 5.2, 2.5, 0.5, size=16, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Problem & Solution
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s3 = add_slide()
bg(s3)
box(s3, 0, 0, 13.33, 0.12, ACCENT_GOLD)
box(s3, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)

section_tag(s3, "PROBLEM & SOLUTION", 0.5, 0.25)
txt(s3, "Challenge & Approach", 0.5, 0.7, 10, 0.7,
    size=32, bold=True, color=ACCENT_GOLD)
divider(s3, 1.5)

# Problem box
box(s3, 0.5, 1.7, 5.7, 3.8, BG_CARD)
box(s3, 0.5, 1.7, 0.08, 3.8, RGBColor(0xE7, 0x4C, 0x3C))
txt(s3, "❌  Problem", 0.75, 1.8, 5.2, 0.5, size=18, bold=True, color=RGBColor(0xE7, 0x4C, 0x3C))
bullet_item(s3, "78 visually distinct tarot cards", 0.65, 2.35, size=15)
bullet_item(s3, "Complex symbols, figures, colors", 0.65, 2.75, size=15)
bullet_item(s3, "Small dataset (~30 photos/card)", 0.65, 3.15, size=15)
bullet_item(s3, "Varied lighting & angles in photos", 0.65, 3.55, size=15)
bullet_item(s3, "Training from scratch = not enough data", 0.65, 3.95, size=15)

# Solution box
box(s3, 6.8, 1.7, 5.9, 3.8, BG_CARD)
box(s3, 6.8, 1.7, 0.08, 3.8, RGBColor(0x2E, 0xCC, 0x71))
txt(s3, "✅  Solution", 7.05, 1.8, 5.4, 0.5, size=18, bold=True, color=RGBColor(0x2E, 0xCC, 0x71))
bullet_item(s3, "Transfer learning from ImageNet", 6.95, 2.35, size=15)
bullet_item(s3, "EfficientNet-B0 as feature extractor", 6.95, 2.75, size=15)
bullet_item(s3, "Data augmentation (flip, rotate, jitter)", 6.95, 3.15, size=15)
bullet_item(s3, "2-phase training (freeze → unfreeze)", 6.95, 3.55, size=15)
bullet_item(s3, "Good accuracy with small dataset", 6.95, 3.95, size=15)

txt(s3, "Key Insight: Borrow knowledge from 1.2M ImageNet images → apply to 78 tarot cards",
    0.5, 5.9, 12.3, 0.6, size=15, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Model Architecture
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s4 = add_slide()
bg(s4)
box(s4, 0, 0, 13.33, 0.12, ACCENT_GOLD)
box(s4, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)

section_tag(s4, "ARCHITECTURE", 0.5, 0.25)
txt(s4, "Model Architecture", 0.5, 0.7, 10, 0.7,
    size=32, bold=True, color=ACCENT_GOLD)
divider(s4, 1.5)

# Architecture flow — horizontal blocks
arch = [
    ("Input\n224×224×3", ACCENT_PURPLE),
    ("EfficientNet-B0\nBackbone\n(frozen Phase 1)", RGBColor(0x2C, 0x3E, 0x7A)),
    ("Global Avg\nPooling\n1280 features", RGBColor(0x2C, 0x3E, 0x7A)),
    ("Dropout\n+Linear\n→512", RGBColor(0x6C, 0x35, 0x8A)),
    ("ReLU\n+Dropout", RGBColor(0x6C, 0x35, 0x8A)),
    ("Linear\n→78\n(Softmax)", ACCENT_GOLD),
]

x = 0.4
for i, (label, color) in enumerate(arch):
    w = 1.9
    box(s4, x, 2.0, w, 1.6, color)
    txt(s4, label, x, 2.1, w, 1.4, size=12, bold=True,
        color=BG_DARK if color == ACCENT_GOLD else TEXT_WHITE,
        align=PP_ALIGN.CENTER)
    if i < len(arch) - 1:
        txt(s4, "→", x + w, 2.55, 0.4, 0.6, size=20, bold=True,
            color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    x += w + 0.4

# Details table
headers = ["Property", "Value"]
rows = [
    ["Parameters",    "~5.3M"],
    ["Input Size",    "224 × 224 px"],
    ["Pretrained On", "ImageNet (1.2M images)"],
    ["Output",        "78 classes (tarot cards)"],
    ["Dropout",       "0.3 (both layers)"],
]

col_w = [3.5, 5.0]
col_x = [0.5, 4.2]
row_h = 0.38
ty = 4.1

for ci, h in enumerate(headers):
    box(s4, col_x[ci], ty, col_w[ci], row_h, ACCENT_PURPLE)
    txt(s4, h, col_x[ci] + 0.1, ty + 0.04, col_w[ci], row_h,
        size=13, bold=True, color=TEXT_WHITE)

for ri, row in enumerate(rows):
    row_color = BG_CARD if ri % 2 == 0 else RGBColor(0x14, 0x14, 0x32)
    for ci, cell in enumerate(row):
        box(s4, col_x[ci], ty + row_h * (ri + 1), col_w[ci], row_h, row_color)
        txt(s4, cell, col_x[ci] + 0.1, ty + row_h * (ri + 1) + 0.04, col_w[ci], row_h,
            size=13, color=TEXT_LIGHT if ci == 0 else TEXT_WHITE)

txt(s4, "Why EfficientNet-B0?  Scales depth + width + resolution together → best accuracy per parameter",
    0.5, 6.9, 12.3, 0.45, size=13, italic=True, color=ACCENT_GOLD)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Training Pipeline
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s5 = add_slide()
bg(s5)
box(s5, 0, 0, 13.33, 0.12, ACCENT_GOLD)
box(s5, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)

section_tag(s5, "TRAINING PIPELINE", 0.5, 0.25)
txt(s5, "Training Pipeline", 0.5, 0.7, 10, 0.7,
    size=32, bold=True, color=ACCENT_GOLD)
divider(s5, 1.5)

# Phase boxes side by side
# Left: Data
box(s5, 0.5, 1.7, 3.7, 4.5, BG_CARD)
box(s5, 0.5, 1.7, 0.08, 4.5, ACCENT_GOLD)
txt(s5, "Data Prep", 0.75, 1.8, 3.3, 0.45, size=15, bold=True, color=ACCENT_GOLD)
for i, t in enumerate(["78 card classes", "~30 photos / card", "HEIC / JPG / PNG", "Split 75/15/10%"]):
    bullet_item(s5, t, 0.65, 2.35 + i * 0.45, w=3.5, size=14)

# Middle: Phase 1
box(s5, 4.5, 1.7, 3.8, 4.5, BG_CARD)
box(s5, 4.5, 1.7, 0.08, 4.5, ACCENT_PURPLE)
txt(s5, "Phase 1  (10 epochs)", 4.75, 1.8, 3.4, 0.45, size=15, bold=True, color=ACCENT_PURPLE)
for i, t in enumerate(["Backbone frozen", "Train head only", "LR = 1e-3", "CosineAnnealingLR", "Fast convergence"]):
    bullet_item(s5, t, 4.65, 2.35 + i * 0.45, w=3.6, size=14)

# Right: Phase 2
box(s5, 8.6, 1.7, 4.2, 4.5, BG_CARD)
box(s5, 8.6, 1.7, 0.08, 4.5, RGBColor(0x2E, 0xCC, 0x71))
txt(s5, "Phase 2  (20 epochs)", 8.85, 1.8, 3.8, 0.45, size=15, bold=True, color=RGBColor(0x2E, 0xCC, 0x71))
for i, t in enumerate(["Backbone unfrozen", "Fine-tune all layers", "LR = 1e-4", "Weight decay 1e-4", "Best model saved"]):
    bullet_item(s5, t, 8.75, 2.35 + i * 0.45, w=3.9, size=14)

txt(s5, "Augmentation: Resize → RandomCrop → Flip → Rotate ±15° → ColorJitter → Normalize",
    0.5, 6.5, 12.3, 0.45, size=13, italic=True, color=ACCENT_GOLD)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 — Roadmap
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s6 = add_slide()
bg(s6)
box(s6, 0, 0, 13.33, 0.12, ACCENT_GOLD)
box(s6, 0, 7.38, 13.33, 0.12, ACCENT_PURPLE)

section_tag(s6, "ROADMAP", 0.5, 0.25)
txt(s6, "Project Roadmap", 0.5, 0.7, 10, 0.7,
    size=32, bold=True, color=ACCENT_GOLD)
divider(s6, 1.5)

phases = [
    ("Phase 1", "Data Collection",    "Photograph all 78 tarot cards\n~30 photos per card, HEIC/JPG",   ACCENT_GOLD,              "🔄 In Progress"),
    ("Phase 2", "Model Training",     "Fine-tune EfficientNet-B0\nGoogle Colab T4 GPU",                  ACCENT_PURPLE,            "⏳ Upcoming"),
    ("Phase 3", "Prediction",         "Identify card from photo\nTop-3 confidence scores",               RGBColor(0x2E, 0xCC, 0x71),"⏳ Upcoming"),
    ("Phase 4", "Fortune Telling",    "Map card → meaning\nGenerate fortune reading",                    RGBColor(0x3B, 0x97, 0xD3),"⏳ Upcoming"),
    ("Phase 5", "App",                "Build full UI\nPhoto capture + fortune display",                  RGBColor(0xE6, 0x7E, 0x22),"⏳ Upcoming"),
]

for i, (phase, title, desc, color, status) in enumerate(phases):
    y = 1.75 + i * 1.0
    box(s6, 0.5, y, 1.5, 0.75, color)
    txt(s6, phase, 0.5, y + 0.15, 1.5, 0.45, size=13, bold=True,
        color=BG_DARK if color == ACCENT_GOLD else TEXT_WHITE, align=PP_ALIGN.CENTER)

    box(s6, 2.2, y, 10.6, 0.75, BG_CARD)
    txt(s6, title, 2.4, y + 0.02, 4.0, 0.38, size=15, bold=True, color=color)
    txt(s6, desc, 2.4, y + 0.38, 7.0, 0.35, size=12, color=TEXT_LIGHT)
    txt(s6, status, 10.0, y + 0.2, 2.6, 0.35, size=13, bold=True,
        color=ACCENT_GOLD if "Progress" in status else TEXT_LIGHT,
        align=PP_ALIGN.RIGHT)


out = "c:/Users/WannakornSangthongng/Desktop/Fortune Teller/Fortune_Teller_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
